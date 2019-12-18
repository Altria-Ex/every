import os
import time
import docx
import pptx
import xlrd
from bs4 import BeautifulSoup
from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdfpage import PDFPage, PDFTextExtractionNotAllowed
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import LAParams
from win32com import client

def disk_search():
    disks = (os.popen('fsutil fsinfo drives').read()).split()[1:]
    disks = disks[1:]
    dirs = list()
    read_files = list()
    #read_ext = {'.txt', '.html', '.pdf', '.c', '.cpp', '.py', '.docx', '.doc', '.pptx', '.ppt', '.xls', '.xlsx'}
    read_ext = {'.html'}
    for disk in disks:
        start_time = time.time()
        dirs.append(list(os.walk(disk)))
        end_time = time.time()
        print(disk, end_time - start_time)
        for root, _, files in dirs[-1]:
            for file in files:
                _, ext = os.path.splitext(file)
                if ext in read_ext:
                   read_files.append(os.path.join(root, file))
    
    return dirs, read_files
    
def build_index(files):




    pass

def read_txt(path):
    text = list()
    if os.path.splitext(path)[1] == '.txt':
        with open (path, 'r') as f:
            for line in f:
                text.append(line)

def read_doc(path):
    text = list()
    if os.path.splitext(path)[1] == '.doc':
        word = client.Dispatch('Word.Application')
        doc = word.Documents.Open(path)
        filepath = os.path.splitext(path)[0] + '.docx'
        doc.SaveAs(filepath, 16)
        doc.Close()
        word.Quit()
        text = read_docx(filepath)
        os.remove(filepath)
    return text

def read_docx(path):     
    text = list()
    if os.path.splitext(path)[1] == '.docx':
        file = docx.Document(path)
        for p in file.paragraphs:
            text.append(p.text)
    return text

def read_ppt(path):
    text = list()
    if os.path.splitext(path)[1] == '.ppt':
        powerpoint = client.Dispatch('PowerPoint.Application')
        ppt = powerpoint.Presentations.Open(path)
        filepath = os.path.splitext(path)[0] + '.pptx'
        ppt.SaveAs(filepath)
        ppt.Close()
        powerpoint.Quit()
        text = read_pptx(filepath)
        os.remove(filepath)
    return text

def read_pptx(path):
    text = list()
    if os.path.splitext(path)[1] == '.pptx':
        ppt = pptx.Presentation(path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text)
    return text


def read_xls_xlsx(path):
    text = list()
    if os.path.splitext(path)[1] == '.xls' or os.path.splitext(path)[1] == '.xlsx':
        workbook = xlrd.open_workbook(path)
        for name in workbook.sheet_names():
            sheet = workbook.sheet_by_name(name)
            for i in range(sheet.nrows):
                temp = sheet.row_values(i)
                text.append(''.join(map(str, temp)).replace(' ', ''))
    return text

def read_pdf(path):
    text = list()
    if os.path.splitext(path)[1] == '.pdf':
        with open(path, 'rb') as f:
            parser = PDFParser(f)
            doc = PDFDocument(parser)
            if doc.is_extractable:
                pdfrm = PDFResourceManager()
                laparams = LAParams()
                device = PDFPageAggregator(pdfrm, laparams=laparams)
                interpreter = PDFPageInterpreter(pdfrm, device)
                for page in PDFPage.create_pages(doc):
                    interpreter.process_page(page)
                    layout = device.get_result()
                    for i in layout:
                        if hasattr(i, 'get_text'):
                            text.append(i.get_text())
        
    return text
        
def read_html(path):
    text = list()
    if os.path.splitext(path)[1] == '.html':
        with open(path, 'r', encoding='utf-8') as htmlfile:
            htmlhandle = htmlfile.read()
            soup = BeautifulSoup(htmlhandle, features='lxml')
            [s.extract() for s in soup(['script', 'style'])]
            text = soup.text.split()
    return text